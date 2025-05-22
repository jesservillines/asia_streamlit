# -*- coding: utf-8 -*-
"""build_presentation.py

Create a self-contained PowerPoint deck with our key visuals.

Slides generated:
1. Title
2. Overview bullets
3. Radar RMSE figure
4. SHAP beeswarm figure
5. Enhanced calibration curve
6. Key findings
7. Next steps
8. Acknowledgements

Run from repo root:
    python -m presentation.build_presentation
Requires `python-pptx`.
"""
from __future__ import annotations

from pathlib import Path
from typing import List, Tuple

from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = ROOT / "asia-impairment-track-prediction" / "visuals" / "figures"
OUTPUT_PATH = Path(__file__).with_name("ASIA_motor_score_presentation.pptx")

# ---------------------------------------------------------------------------
# Slide layout indices (default template)
# ---------------------------------------------------------------------------
TITLE_SLIDE = 0
TITLE_AND_CONTENT = 1
BLANK = 6

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def _set_font(paragraph, size: int) -> None:
    paragraph.font.size = Pt(size)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE])
    slide.shapes.title.text = (
        "Predicting Long-Term Motor Recovery after Spinal Cord Injury"
    )
    subtitle = slide.placeholders[1]
    subtitle.text = "Craig Hospital • ASIA 2025 • Presenter: <Name>"


def add_bullet_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_AND_CONTENT])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        _set_font(p, 18)
        p.level = 0


def add_image_slide(prs: Presentation, title: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])

    # Title
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = tbox.text_frame
    tf.text = title
    tf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    _set_font(tf.paragraphs[0], 24)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    if image_path.exists():
        pic = slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(8))
        pic.left = int((prs.slide_width - pic.width) / 2)
    else:
        add_bullet_slide(prs, title, [f"[Missing figure: {image_path.name}]"])

# ---------------------------------------------------------------------------
# Main build
# ---------------------------------------------------------------------------

def main() -> None:
    prs = Presentation()

    # 1. Title slide
    add_title_slide(prs)

    # 2. Overview
    add_bullet_slide(
        prs,
        "Project Overview",
        [
            "Objective: predict 20 limb motor scores at weeks 26 & 52 using week-1 exam data.",
            "Dataset: 544 patients, ~150 baseline clinical features.",
            "Model: Ensemble of CatBoost, XGBoost, HistGradientBoosting.",
            "Evaluation: macro-averaged RMSE (competition metric).",
        ],
    )

    # New slide: Data processing / feature engineering
    add_bullet_slide(
        prs,
        "Data Processing & Feature Engineering",
        [
            "Dropped PID identifier; imputed missing values (median/mode).",
            "One-hot encoded AIS grade, neurological level, and categorical vars.",
            "Derived composite strength scores & age bands.",
            "Scaled continuous predictors (z-score) prior to modeling.",
        ],
    )

    # New slide: Modelling methodology
    add_bullet_slide(
        prs,
        "Modelling Methodology",
        [
            "Base models: CatBoost, XGBoost, HistGradientBoosting (20 outputs each).",
            "Hyperparameters tuned via Optuna with 5-fold CV, early stopping.",
            "Final ensemble = equal-weighted mean of per-target predictions.",
            "Validation: patient-level 80/20 split to avoid leakage.",
        ],
    )

    # 3-5. Figures
    figures: List[Tuple[str, str]] = [
        ("Per-Target RMSE (Radar)", "radar_target_rmse.png"),
        ("Feature Importance (SHAP Beeswarm)", "shap_beeswarm_ensemble.png"),
        ("Calibration Curve", "calibration_curve_enhanced.png"),
        ("Residual Heatmap", "residuals_heatmap.png"),
        ("Best vs Worst Patient Feature Diff", "patient_group_diff_bars.png"),
        ("Demographic Composition", "patient_group_demographics.png"),
    ]
    for title, fname in figures:
        add_image_slide(prs, title, FIG_DIR / fname)

    # 6. Findings slide
    add_bullet_slide(
        prs,
        "Key Findings",
        [
            "Ensemble attains RMSE < 0.90 on validation set.",
            "Predictions well-calibrated at 0/5 scores; slight over-prediction mid-range.",
            "Age, initial AIS grade, proximal strength most influential features.",
        ],
    )

    # 7. Next steps
    add_bullet_slide(
        prs,
        "Next Steps",
        [
            "Incorporate MRI imaging & longitudinal fine-tuning.",
            "Develop bedside decision-support prototype.",
        ],
    )

    # 8. Acknowledgements
    add_bullet_slide(
        prs,
        "Acknowledgements",
        [
            "Craig Hospital Research Department & SCI Model Systems.",
            "Collaborators, clinicians, and patients contributing data.",
            "Kaggle & ASIA organising committee.",
        ],
    )

    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to {OUTPUT_PATH.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
